"""
PPTX Ingestor
Extracts text from PowerPoint presentations
"""

import logging
from typing import List, Dict
from pathlib import Path
from .chunker import document_chunker

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

logger = logging.getLogger(__name__)

def ingest_pptx(
    file_path: str,
    file_id: str,
    progress_callback=None
) -> List[Dict]:
    """
    Ingest a PPTX file and return chunks.
    
    Args:
        file_path: Path to the .pptx file
        file_id: Unique ID for the file
        progress_callback: Optional callable(step_number, message, percentage)
    """
    try:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {path}")

        if progress_callback:
            progress_callback(1, "Loading PowerPoint file...", 10)

        if not HAS_PPTX:
            raise ImportError("python-pptx library not found. Please install it with 'pip install python-pptx'")

        prs = Presentation(str(path))
        pages = []
        total_slides = len(prs.slides)
        
        logger.info(f"Extracting text from {total_slides} slides in {path.name}")

        for i, slide in enumerate(prs.slides):
            slide_text_parts = []
            
            # Sort shapes by position (top-to-bottom, left-to-right) for better reading order
            # shape.top and shape.left exist on most shapes
            shapes = []
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    shapes.append(shape)
            
            # Simple sorting heuristic: primarily by top, then by left
            shapes.sort(key=lambda s: (s.top, s.left))

            for shape in shapes:
                text = shape.text
                if text and text.strip():
                    slide_text_parts.append(text.strip())

            full_slide_text = "\n\n".join(slide_text_parts)
            
            if full_slide_text.strip():
                pages.append({
                    "page": i + 1,
                    "text": full_slide_text,
                    "file_id": file_id
                })
            
            if progress_callback:
                # Progress from 10% to 50% during extraction
                percent = 10 + int((i / total_slides) * 40)
                progress_callback(1, f"Extracting slide {i+1}/{total_slides}", percent)

        if not pages:
            logger.warning(f"No text content found in {path.name}")
            return []

        # Chunking (progress handled inside chunker, mapping to step 2)
        if progress_callback:
            progress_callback(2, "Chunking content...", 50)
            
        chunks = document_chunker.chunk_document_pages(
            pages,
            progress_file_id=None # We are manually handling progress via callback if needed, but chunker supports file_id
        )
        
        # document_chunker doesn't take the generic callback, it takes progress_file_id.
        # But in qa_service we wrapped the callback. 
        # Actually document_chunker.chunk_document_pages takes progress_file_id (string) to call progress_service directly.
        # Since we are inside ingestion, we might just return chunks. The qa_service calls us.
        # Let's check qa_service logic again. 
        # qa_service passes 'progress_file_id' to chunker. 
        # But for other ingestors (like excel) it uses 'progress_callback'. 
        # Let's standardize on returning chunks.

        return chunks

    except Exception as e:
        logger.error(f"Error ingesting PPTX {file_path}: {e}")
        raise
